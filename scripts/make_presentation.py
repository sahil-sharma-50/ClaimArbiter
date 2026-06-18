"""Generate the ClaimArbiter slide deck (16:9 PPTX) into assets/.

Design language follows the reference deck: near-black canvas, a single warm
accent, a top-left logo lockup on every slide, an accent kicker label, a big
bold white headline, dark rounded cards with outline-pill tags, a pipeline row
with arrows, and a "powered by" sponsor footer. Retitled and re-themed for
ClaimArbiter (gold accent + the real gold robot brandmark).

Five slides: Cover, Problem, Solution, Technology + sponsors, Live demo CTA.

Run (python-pptx is pulled on demand, no repo dependency added):

    uv run --with python-pptx python scripts/make_presentation.py

Writes assets/ClaimArbiter-Presentation.pptx.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

# ── 16:9 canvas (13.333in x 7.5in) ──
EMU_IN = 914400
SW = int(13.333 * EMU_IN)
SH = int(7.5 * EMU_IN)

# ── Aurora Noir palette (echoes frontend/app/globals.css) ──
CANVAS = RGBColor(0x0B, 0x0D, 0x14)      # near-black navy
PANEL = RGBColor(0x15, 0x18, 0x22)       # raised card
PANEL_HI = RGBColor(0x1B, 0x1F, 0x2B)    # elevated
INK = RGBColor(0xF2, 0xF3, 0xF6)         # primary text
SOFT = RGBColor(0xB7, 0xBB, 0xC6)        # secondary text
FAINT = RGBColor(0x8B, 0x90, 0x9C)       # faint
GHOST = RGBColor(0x66, 0x6B, 0x78)       # footer
GOLD = RGBColor(0xE8, 0xB5, 0x4A)        # signature accent
GOLD_DK = RGBColor(0xC4, 0xA0, 0x35)
LINE = RGBColor(0x2B, 0x2E, 0x39)        # hairline
ORG_A = RGBColor(0x86, 0x92, 0xDC)       # insurer (blue-violet)
ORG_B = RGBColor(0xD4, 0xA7, 0x73)       # specialist (copper)
SUCCESS = RGBColor(0x4C, 0xCB, 0x8D)

FONT = "Arial"  # universal; matches the clean grotesk feel of the source deck

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "assets" / "ClaimArbiter-Presentation.pptx"
LOGO = ROOT / "assets" / "logo-mark.png"


def _in(v: float) -> int:
    return int(v * EMU_IN)


def solid(shape, color) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_line(shape) -> None:
    shape.line.fill.background()


def line(shape, color, w: float = 1.0) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(w)


def no_shadow(shape) -> None:
    # python-pptx has no shadow toggle; strip the inherited preset via XML.
    sp = shape._element.spPr
    existing = sp.find(qn("a:effectLst"))
    if existing is None:
        sp.append(sp.makeelement(qn("a:effectLst"), {}))


def textbox(slide, x, y, w, h, *, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(_in(x), _in(y), _in(w), _in(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def set_run(run, text, size, color, *, bold=False, italic=False,
            spacing=None, font=FONT):
    run.text = text
    f = run.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    if spacing is not None:
        # letter spacing in 1/100 pt
        run.font._rPr.set("spc", str(int(spacing * 100)))


def para(tf, *, first=False, space_after=0.0, space_before=0.0, line_spacing=None,
         align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if space_after:
        p.space_after = Pt(space_after)
    if space_before:
        p.space_before = Pt(space_before)
    if line_spacing:
        p.line_spacing = line_spacing
    return p


def bg(slide) -> None:
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    solid(r, CANVAS)
    no_line(r)
    no_shadow(r)


def gold_spine(slide) -> None:
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, _in(0.07), SH)
    solid(r, GOLD)
    no_line(r)
    no_shadow(r)


def logo_lockup(slide, *, small=True) -> None:
    """Top-left brandmark + ClaimArbiter wordmark, as in the reference deck."""
    s = 0.42 if small else 0.62
    yld = 0.42
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), _in(0.55), _in(yld), _in(s), _in(s))
    tb, tf = textbox(slide, 0.55 + s + 0.14, yld - 0.04, 4.5, s + 0.1,
                     anchor=MSO_ANCHOR.MIDDLE)
    p = para(tf, first=True)
    r1 = p.add_run(); set_run(r1, "Claim", 16 if small else 20, INK, bold=True)
    r2 = p.add_run(); set_run(r2, "Arbiter", 16 if small else 20, GOLD, bold=True)


def kicker(slide, x, yy, text, color=GOLD) -> None:
    tb, tf = textbox(slide, x, yy, 8, 0.32)
    p = para(tf, first=True)
    r = p.add_run()
    set_run(r, text.upper(), 11.5, color, bold=True, spacing=2.4)


def powered_by(slide, sponsors=("Band", "AI/ML API", "Featherless")) -> None:
    """'powered by' + sponsor wordmarks, bottom-right (sponsor credit)."""
    tb, tf = textbox(slide, 6.4, 6.86, 6.4, 0.4, anchor=MSO_ANCHOR.MIDDLE)
    p = para(tf, first=True, align=PP_ALIGN.RIGHT)
    r = p.add_run(); set_run(r, "powered by   ", 11, GHOST)
    for i, s in enumerate(sponsors):
        rr = p.add_run()
        set_run(rr, s, 12.5, INK, bold=True)
        if i < len(sponsors) - 1:
            sep = p.add_run(); set_run(sep, "   ·   ", 12.5, GOLD)


def card(slide, x, y, w, h, *, fill=PANEL, border=LINE, border_w=1.0):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, _in(x), _in(y),
                               _in(w), _in(h))
    c.adjustments[0] = 0.06
    solid(c, fill)
    line(c, border, border_w)
    no_shadow(c)
    return c


def pill(slide, x, y, text, color=GOLD):
    """Outline pill tag, like CAUSE / EFFECT / COST in the reference."""
    w = 0.30 + 0.092 * len(text)
    h = 0.30
    p = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, _in(x), _in(y),
                               _in(w), _in(h))
    p.adjustments[0] = 0.5
    p.fill.background()
    line(p, color, 1.25)
    no_shadow(p)
    tf = p.text_frame
    tf.word_wrap = False
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    para_ = tf.paragraphs[0]; para_.alignment = PP_ALIGN.CENTER
    r = para_.add_run(); set_run(r, text.upper(), 8.5, color, bold=True, spacing=1.2)
    return w


def accent_bar(slide, x, y, color=GOLD, w=0.5, h=0.045):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _in(x), _in(y), _in(w), _in(h))
    solid(r, color)
    no_line(r)
    no_shadow(r)


def arrow(slide, x, y, w=0.55, color=GOLD):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, _in(x), _in(y), _in(w), _in(0.22))
    solid(a, color)
    no_line(a)
    no_shadow(a)
    a.adjustments[0] = 0.55
    a.adjustments[1] = 0.55


# ───────────────────────────── slides ─────────────────────────────

def slide_cover(prs) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    # centered logo + wordmark
    logo_in = 1.0
    s.shapes.add_picture(str(LOGO), _in(4.62), _in(2.18), _in(logo_in), _in(logo_in))
    tb, tf = textbox(s, 5.7, 2.16, 6.8, 1.1, anchor=MSO_ANCHOR.MIDDLE)
    p = para(tf, first=True)
    r1 = p.add_run(); set_run(r1, "Claim", 60, INK, bold=True)
    r2 = p.add_run(); set_run(r2, "Arbiter", 60, GOLD, bold=True)
    # accent tick
    accent_bar(s, 6.17, 3.66, GOLD, w=1.0, h=0.06)
    # tagline
    tb, tf = textbox(s, 2.0, 3.95, 9.33, 0.6)
    p = para(tf, first=True, align=PP_ALIGN.CENTER)
    r = p.add_run()
    set_run(r, "One claim in, the right specialist recruited, a signed verdict out.",
            20, SOFT, bold=True)
    # event line
    tb, tf = textbox(s, 2.0, 4.55, 9.33, 0.4)
    p = para(tf, first=True, align=PP_ALIGN.CENTER)
    r = p.add_run()
    set_run(r, "Cross-org, cross-framework claim adjudication on Band.", 13.5, FAINT)
    powered_by(s)


def slide_problem(prs) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s); gold_spine(s); logo_lockup(s)
    kicker(s, 0.55, 1.45, "The problem")
    tb, tf = textbox(s, 0.52, 1.78, 12.3, 1.0)
    p = para(tf, first=True)
    r = p.add_run()
    set_run(r, "Insurers do not employ the specialists they depend on.", 33, INK, bold=True)
    tb, tf = textbox(s, 0.55, 2.62, 11.8, 0.5)
    p = para(tf, first=True)
    r = p.add_run()
    set_run(r, "Every property, injury, or legal claim waits on an outside expert reached by email.",
            14.5, FAINT)

    cards = [
        ("CAUSE", "Manual handoffs",
         "Claims sit in a queue, then get emailed to an outside assessor with no shared context."),
        ("EFFECT", "Slow and inconsistent",
         "Decisions return days later. Different reviewers reach different calls on similar claims."),
        ("COST", "Weak audit trail",
         "Who decided what, and why, is reconstructed after the fact. Disputes and compliance suffer."),
    ]
    cx, cy, cw, ch, gap = 0.55, 3.4, 3.94, 2.7, 0.27
    for i, (tag, head, body) in enumerate(cards):
        x = cx + i * (cw + gap)
        card(s, x, cy, cw, ch)
        accent_bar(s, x, cy + 0.0, GOLD, w=0.06, h=ch)  # left gold edge
        pill(s, x + 0.32, cy + 0.34, tag)
        tb, tf = textbox(s, x + 0.32, cy + 0.84, cw - 0.62, 0.5)
        p = para(tf, first=True)
        r = p.add_run(); set_run(r, head, 17, INK, bold=True)
        tb, tf = textbox(s, x + 0.32, cy + 1.38, cw - 0.62, 1.2)
        p = para(tf, first=True, line_spacing=1.18)
        r = p.add_run(); set_run(r, body, 12.5, SOFT)
    powered_by(s)


def slide_solution(prs) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s); gold_spine(s); logo_lockup(s)
    kicker(s, 0.55, 1.45, "The solution")
    tb, tf = textbox(s, 0.52, 1.78, 12.4, 1.0)
    p = para(tf, first=True)
    r1 = p.add_run(); set_run(r1, "One claim ", 33, INK, bold=True)
    r2 = p.add_run(); set_run(r2, "→", 33, GOLD, bold=True)
    r3 = p.add_run(); set_run(r3, " the right specialist, a signed verdict.", 33, INK, bold=True)

    # pipeline row: claim in -> Band agents -> signed verdict
    py = 2.95
    b1 = card(s, 0.55, py, 2.7, 0.92, fill=PANEL_HI)
    tb = b1.text_frame; tb.word_wrap = True
    tb.vertical_anchor = MSO_ANCHOR.MIDDLE
    tb.margin_left = _in(0.1); tb.margin_right = _in(0.1)
    pp = tb.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
    r = pp.add_run(); set_run(r, "Claim filed", 14, INK, bold=True)
    pp2 = tb.add_paragraph(); pp2.alignment = PP_ALIGN.CENTER
    r = pp2.add_run(); set_run(r, "property · medical · legal", 10.5, FAINT)
    arrow(s, 3.45, py + 0.33)

    b2 = card(s, 4.25, py, 4.65, 0.92, fill=PANEL_HI, border=GOLD, border_w=1.25)
    tb = b2.text_frame; tb.word_wrap = True
    tb.vertical_anchor = MSO_ANCHOR.MIDDLE
    tb.margin_left = _in(0.1); tb.margin_right = _in(0.1)
    pp = tb.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
    r = pp.add_run(); set_run(r, "Agents collaborate on Band", 14, GOLD, bold=True)
    pp2 = tb.add_paragraph(); pp2.alignment = PP_ALIGN.CENTER
    r = pp2.add_run(); set_run(r, "intake · evidence · coordinator · recruit · relay", 10.5, SOFT)
    arrow(s, 9.1, py + 0.33)

    b3 = card(s, 9.9, py, 2.88, 0.92, fill=PANEL_HI)
    tb = b3.text_frame; tb.word_wrap = True
    tb.vertical_anchor = MSO_ANCHOR.MIDDLE
    tb.margin_left = _in(0.1); tb.margin_right = _in(0.1)
    pp = tb.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
    r = pp.add_run(); set_run(r, "Signed verdict", 14, SUCCESS, bold=True)
    pp2 = tb.add_paragraph(); pp2.alignment = PP_ALIGN.CENTER
    r = pp2.add_run(); set_run(r, "human approves · audit sealed", 10.5, FAINT)

    # three product cards (domains)
    cards = [
        ("PROPERTY", ORG_B, "Property Group",
         "Assesses water and structural damage against the covered peril."),
        ("MEDICAL", ORG_B, "Medical Group",
         "Checks billed treatment against the reported injury."),
        ("LEGAL", ORG_B, "Legal Group",
         "Reviews lawyer fees and proceedings, and writes the deny rationale."),
    ]
    cx, cy, cw, ch, gap = 0.55, 4.25, 3.94, 2.05, 0.27
    for i, (tag, accent, head, body) in enumerate(cards):
        x = cx + i * (cw + gap)
        card(s, x, cy, cw, ch)
        accent_bar(s, x + 0.32, cy + 0.34, accent, w=0.34, h=0.05)
        pill(s, x + 0.32, cy + 0.46, tag, color=accent)
        tb, tf = textbox(s, x + 0.32, cy + 0.92, cw - 0.64, 0.4)
        p = para(tf, first=True)
        r = p.add_run(); set_run(r, head, 15, INK, bold=True)
        tb, tf = textbox(s, x + 0.32, cy + 1.32, cw - 0.64, 0.7)
        p = para(tf, first=True, line_spacing=1.12)
        r = p.add_run(); set_run(r, body, 11.5, SOFT)
    powered_by(s)


def slide_tech(prs) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s); gold_spine(s); logo_lockup(s)
    kicker(s, 0.55, 1.45, "How it is built")
    tb, tf = textbox(s, 0.52, 1.78, 12.4, 1.0)
    p = para(tf, first=True)
    r = p.add_run()
    set_run(r, "Genuinely cross-framework, genuinely cross-org.", 33, INK, bold=True)

    cards = [
        ("FRAMEWORKS", "Three, one room",
         "Pydantic AI, LangGraph, and CrewAI agents collaborate through Band adapters."),
        ("MODELS", "Frontier + open",
         "GPT-4o at the insurer (AI/ML API); Llama 3.1 at the specialists (Featherless)."),
        ("RECORD", "Band is the truth",
         "Messages, task state, and a tamper-evident audit trail live in Band, not any one agent."),
    ]
    cx, cy, cw, ch, gap = 0.55, 2.95, 3.94, 2.5, 0.27
    for i, (tag, head, body) in enumerate(cards):
        x = cx + i * (cw + gap)
        card(s, x, cy, cw, ch)
        pill(s, x + 0.32, cy + 0.32, tag)
        tb, tf = textbox(s, x + 0.32, cy + 0.82, cw - 0.64, 0.5)
        p = para(tf, first=True)
        r = p.add_run(); set_run(r, head, 17, INK, bold=True)
        tb, tf = textbox(s, x + 0.32, cy + 1.36, cw - 0.64, 1.0)
        p = para(tf, first=True, line_spacing=1.18)
        r = p.add_run(); set_run(r, body, 12.5, SOFT)

    # sponsor strip
    tb, tf = textbox(s, 0.55, 5.75, 12.2, 0.9, anchor=MSO_ANCHOR.MIDDLE)
    p = para(tf, first=True)
    r = p.add_run(); set_run(r, "Built for the Band of Agents Hackathon, powered by  ", 13, FAINT)
    for i, (name, col) in enumerate([("Band", GOLD), ("AI/ML API", ORG_A), ("Featherless", ORG_B)]):
        rr = p.add_run(); set_run(rr, name, 15, col, bold=True)
        if i < 2:
            sep = p.add_run(); set_run(sep, "   ·   ", 14, GHOST)
    powered_by(s)


def slide_demo(prs) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s); gold_spine(s); logo_lockup(s)
    # centered CTA pill, like "Now let's see it run."
    bar = slide_pill = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, _in(0.95), _in(3.05), _in(11.43), _in(1.4))
    slide_pill.adjustments[0] = 0.5
    slide_pill.fill.background()
    line(slide_pill, GOLD, 1.5)
    no_shadow(slide_pill)
    tf = slide_pill.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
    r = pp.add_run(); set_run(r, "Now let's see it run.", 30, GOLD, bold=True)
    # drawn play triangle (glyph-free, renders identically everywhere),
    # placed just left of the centered headline and vertically on the pill axis.
    tri = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, _in(4.12), _in(3.57),
                             _in(0.32), _in(0.34))
    tri.rotation = 90
    solid(tri, GOLD); no_line(tri); no_shadow(tri)

    tb, tf = textbox(s, 0.95, 4.65, 11.43, 0.5)
    p = para(tf, first=True, align=PP_ALIGN.CENTER)
    r = p.add_run()
    set_run(r, "localhost:3000  →  Open platform  →  Live  →  pick property, medical, or legal.",
            14, FAINT)
    powered_by(s)


def build() -> None:
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    slide_cover(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_tech(prs)
    slide_demo(prs)
    prs.save(str(OUT))
    print(f"Wrote {OUT.relative_to(ROOT)} ({len(prs.slides._sldIdLst)} slides, 16:9)")


if __name__ == "__main__":
    build()
