"""Insert a 'Business value & market' slide into the EXISTING deck, in place.

Non-destructive: opens assets/ClaimArbiter-Presentation.pptx (which may carry
manual PowerPoint edits), appends one new slide built with the same styling
helpers as make_presentation.py, then reorders it to sit just before the final
'Now let's see it run' CTA. Slides authored elsewhere are left untouched.

    uv run --with python-pptx python scripts/add_business_slide.py
"""

from __future__ import annotations

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

import make_presentation as M  # reuse palette + styling helpers (no build() runs on import)

OUT = M.OUT


def slide_business(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    M.bg(s); M.gold_spine(s); M.logo_lockup(s)
    M.kicker(s, 0.55, 1.45, "Business value")
    tb, tf = M.textbox(s, 0.52, 1.78, 12.4, 1.0)
    p = M.para(tf, first=True)
    r = p.add_run()
    M.set_run(r, "A large, measurable market.", 33, M.INK, bold=True)

    cards = [
        ("TAM", "Claims operations spend",
         "Loss-adjustment expense runs an estimated 10 to 12 percent of global P&C premiums, a multi-hundred-billion-dollar annual pool."),
        ("SAM", "Digital-first carriers",
         "P&C and health insurers and InsurTechs adopting agentic automation, an estimated multi-billion-dollar serviceable slice."),
        ("WEDGE", "Cross-org adjudication",
         "Claims that need an outside specialist: the slow, manual, high-value handoffs ClaimArbiter automates first."),
    ]
    cx, cy, cw, ch, gap = 0.55, 2.95, 3.94, 2.5, 0.27
    for i, (tag, head, body) in enumerate(cards):
        x = cx + i * (cw + gap)
        M.card(s, x, cy, cw, ch)
        M.pill(s, x + 0.32, cy + 0.32, tag)
        tb, tf = M.textbox(s, x + 0.32, cy + 0.82, cw - 0.64, 0.5)
        p = M.para(tf, first=True)
        r = p.add_run(); M.set_run(r, head, 16, M.INK, bold=True)
        tb, tf = M.textbox(s, x + 0.32, cy + 1.3, cw - 0.64, 1.1)
        p = M.para(tf, first=True, line_spacing=1.18)
        r = p.add_run(); M.set_run(r, body, 12, M.SOFT)

    # revenue strip (mirrors the sponsor strip on the tech slide)
    tb, tf = M.textbox(s, 0.55, 5.72, 12.2, 0.5, anchor=MSO_ANCHOR.MIDDLE)
    p = M.para(tf, first=True)
    r = p.add_run(); M.set_run(r, "Revenue:  ", 14, M.FAINT, bold=True)
    streams = ["Platform subscription", "Usage pricing", "Recruiting take-rate", "Compliance add-on"]
    for i, name in enumerate(streams):
        rr = p.add_run(); M.set_run(rr, name, 14, M.INK, bold=True)
        if i < len(streams) - 1:
            sep = p.add_run(); M.set_run(sep, "   ·   ", 14, M.GOLD)

    # framing disclaimer (honest about the figures)
    tb, tf = M.textbox(s, 0.55, 6.2, 12.2, 0.35)
    p = M.para(tf, first=True)
    r = p.add_run()
    M.set_run(r, "Figures are order-of-magnitude estimates for framing, not audited market data.",
              10.5, M.GHOST)
    M.powered_by(s)
    return s


def move_before_last(prs) -> None:
    """Move the just-added slide (currently last) to second-from-last, so it
    lands immediately before the closing CTA slide."""
    sld_lst = prs.slides._sldIdLst
    ids = list(sld_lst)
    if len(ids) < 2:
        return
    new = ids[-1]            # the slide we just appended
    sld_lst.remove(new)
    sld_lst.insert(len(ids) - 2, new)  # before the old last (the CTA)


def main() -> None:
    prs = Presentation(str(OUT))
    before = len(prs.slides._sldIdLst)
    slide_business(prs)
    move_before_last(prs)
    prs.save(str(OUT))
    after = len(prs.slides._sldIdLst)
    print(f"Inserted Business value slide. {before} -> {after} slides; "
          f"it now sits at position {after - 1} of {after} (before the CTA).")


if __name__ == "__main__":
    main()
